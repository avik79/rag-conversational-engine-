"""Handles ingestion and text extraction for all supported document formats."""

import io
from pathlib import Path
from typing import List

import pandas as pd
import pypdf
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook


SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".xlsx", ".xls", ".txt", ".csv", ".md"}


def extract_text(file_name: str, file_bytes: bytes) -> List[dict]:
    """
    Extract text chunks from a file.
    Returns list of {"text": str, "source": str, "page": int} dicts.
    """
    ext = Path(file_name).suffix.lower()
    if ext == ".pdf":
        return _extract_pdf(file_name, file_bytes)
    elif ext == ".docx":
        return _extract_docx(file_name, file_bytes)
    elif ext == ".pptx":
        return _extract_pptx(file_name, file_bytes)
    elif ext in (".xlsx", ".xls"):
        return _extract_excel(file_name, file_bytes)
    elif ext == ".csv":
        return _extract_csv(file_name, file_bytes)
    elif ext in (".txt", ".md"):
        return _extract_text(file_name, file_bytes)
    else:
        raise ValueError(f"Unsupported file type: {ext}")


def _extract_pdf(name: str, data: bytes) -> List[dict]:
    reader = pypdf.PdfReader(io.BytesIO(data))
    chunks = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text() or ""
        text = text.strip()
        if text:
            chunks.append({"text": text, "source": name, "page": i + 1})
    return chunks


def _extract_docx(name: str, data: bytes) -> List[dict]:
    doc = Document(io.BytesIO(data))
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    # Group paragraphs into ~500-word chunks
    return _chunk_paragraphs(paragraphs, name)


def _extract_pptx(name: str, data: bytes) -> List[dict]:
    prs = Presentation(io.BytesIO(data))
    chunks = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
        if texts:
            chunks.append({"text": "\n".join(texts), "source": name, "page": i + 1})
    return chunks


def _extract_excel(name: str, data: bytes) -> List[dict]:
    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    chunks = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            row_text = "\t".join(str(c) if c is not None else "" for c in row)
            if row_text.strip():
                rows.append(row_text)
        if rows:
            chunks.append({
                "text": "\n".join(rows),
                "source": name,
                "page": sheet_name,
            })
    return chunks


def _extract_csv(name: str, data: bytes) -> List[dict]:
    df = pd.read_csv(io.BytesIO(data))
    text = df.to_string(index=False)
    return [{"text": text, "source": name, "page": 1}]


def _extract_text(name: str, data: bytes) -> List[dict]:
    text = data.decode("utf-8", errors="replace").strip()
    # Split into ~1000-char chunks on paragraph boundaries
    paragraphs = [p.strip() for p in text.split("\n\n") if p.strip()]
    return _chunk_paragraphs(paragraphs, name)


def _chunk_paragraphs(paragraphs: List[str], source: str, max_words: int = 300) -> List[dict]:
    chunks, current, word_count, page = [], [], 0, 1
    for para in paragraphs:
        words = len(para.split())
        if word_count + words > max_words and current:
            chunks.append({"text": " ".join(current), "source": source, "page": page})
            current, word_count = [], 0
            page += 1
        current.append(para)
        word_count += words
    if current:
        chunks.append({"text": " ".join(current), "source": source, "page": page})
    return chunks
